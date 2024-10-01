import collections 
import collections.abc
from pptx import Presentation

path = r'C:\Users\Aymen Mahmoudi\Lenovo_Sync\Programming\Python Projects\PPT Notes Extractor\tests'
file = 'test_pptx_file.pptx'

ppt=Presentation(path + '\\'+file)

notes = []

for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note. 
    textNote = slide.notes_slide.notes_text_frame.text
    notes.append((page,textNote)) 
    

#print(notes)

#print(type(notes[1][1]))

# Open a text file in write mode ('w')
with open(path+ '\\'+file[:-4]+'_Notes.txt', 'w', encoding='utf-8') as f:
    for page, slide in enumerate(ppt.slides):
        # this is the notes that doesn't appear on the ppt slide,
        # but really the 'presenter' note. 
        textNote = slide.notes_slide.notes_text_frame.text
        # Write some text to the file
        f.write(str(page + 1) + "\n")
        f.write('================================'+ "\n")
        f.write(textNote+ "\n")




print('Done !!')